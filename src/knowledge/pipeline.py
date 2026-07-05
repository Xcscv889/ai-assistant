"""文档摄入流水线"""

import os
from pathlib import Path
from typing import Any, Dict, List, Optional

from loguru import logger

from .chunker import DocumentChunker
from .embedding import EmbeddingService
from .vector_store import VectorStore


class DocumentPipeline:
    """文档摄入流水线：解析 → 分块 → Embedding → 入库"""

    def __init__(
        self,
        embedding_service: Optional[EmbeddingService] = None,
        vector_store: Optional[VectorStore] = None,
        chunker: Optional[DocumentChunker] = None,
    ):
        self.embedding_service = embedding_service or EmbeddingService(provider="auto")
        self.vector_store = vector_store or VectorStore()
        self.chunker = chunker or DocumentChunker()

    def ingest_file(self, file_path: str, metadata: Optional[Dict[str, Any]] = None) -> List[str]:
        """摄入单个文件到知识库

        Args:
            file_path: 文件路径
            metadata: 额外的元数据

        Returns:
            添加的文档 ID 列表
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        # 解析文件内容
        text_content = self._parse_file(path)

        if not text_content:
            logger.warning(f"文件内容为空: {file_path}")
            return []

        # 构建元数据
        meta = {
            "source": str(path.absolute()),
            "filename": path.name,
            "file_type": path.suffix.lower(),
            "file_size_mb": round(path.stat().st_size / (1024 * 1024), 2),
            **(metadata or {}),
        }

        # 分块
        chunks = self.chunker.chunk(text_content, metadata=meta)

        if not chunks:
            return []

        # 生成 Embedding
        texts = [c["content"] for c in chunks]
        embeddings = self.embedding_service.embed(texts)

        # 入库
        ids = self.vector_store.add(
            documents=texts,
            embeddings=embeddings,
            metadatas=[c["metadata"] for c in chunks],
        )

        logger.info(f"文档已入库: {path.name} ({len(chunks)} 个块)")
        return ids

    def ingest_directory(
        self, directory: str, recursive: bool = True, metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, List[str]]:
        """摄入目录中的所有文件

        Returns:
            {文件路径: [文档ID列表]}
        """
        from ..utils.file_utils import scan_directory

        results = {}
        files = scan_directory(Path(directory), recursive=recursive)

        for file_path in files:
            try:
                ids = self.ingest_file(str(file_path), metadata=metadata)
                results[str(file_path)] = ids
            except Exception as e:
                logger.error(f"摄入文件失败 {file_path}: {e}")
                results[str(file_path)] = []

        total_chunks = sum(len(ids) for ids in results.values())
        logger.info(f"目录摄入完成: {len(files)} 个文件, {total_chunks} 个块")
        return results

    def search(
        self, query: str, top_k: Optional[int] = None
    ) -> List[Dict[str, Any]]:
        """搜索知识库

        Args:
            query: 查询文本
            top_k: 返回结果数量

        Returns:
            搜索结果
        """
        query_embedding = self.embedding_service.embed_query(query)
        return self.vector_store.search(query_embedding, top_k=top_k)

    def _parse_file(self, file_path: Path) -> str:
        """根据文件类型选择合适的解析器"""
        suffix = file_path.suffix.lower()

        if suffix in (".txt", ".md"):
            return self._parse_text(file_path)
        elif suffix == ".pdf":
            return self._parse_pdf(file_path)
        elif suffix == ".docx":
            return self._parse_docx(file_path)
        elif suffix == ".pptx":
            return self._parse_pptx(file_path)
        elif suffix == ".xlsx":
            return self._parse_xlsx(file_path)
        else:
            raise ValueError(f"不支持的文件格式: {suffix}")

    def _parse_text(self, file_path: Path) -> str:
        """解析纯文本文件"""
        from ..utils.file_utils import read_text_file

        content = read_text_file(file_path)
        return content or ""

    def _parse_pdf(self, file_path: Path) -> str:
        """解析 PDF 文件"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(file_path))
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()
            return "\n\n".join(text_parts)
        except ImportError:
            raise ImportError("请安装 pymupdf: pip install pymupdf")

    def _parse_docx(self, file_path: Path) -> str:
        """解析 Word 文档"""
        try:
            from docx import Document

            doc = Document(str(file_path))
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 也解析表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n\n".join(text_parts)
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

    def _parse_pptx(self, file_path: Path) -> str:
        """解析 PPT 文档"""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text)

                    # 解析表格
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                            if row_text.strip():
                                slide_texts.append(row_text)

                text_parts.append("\n".join(slide_texts))

            return "\n\n".join(text_parts)
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

    def _parse_xlsx(self, file_path: Path) -> str:
        """解析 Excel 文件"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_texts = [f"--- Sheet: {sheet_name} ---"]

                for row in ws.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        sheet_texts.append(" | ".join(row_values))

                text_parts.append("\n".join(sheet_texts))

            wb.close()
            return "\n\n".join(text_parts)
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")
