"""MCP文档解析工具"""

from pathlib import Path
from typing import Any, Dict, Optional

from loguru import logger


class DocumentParserTool:
    """文档解析MCP工具 - 支持PDF/Word/PPT/Excel/Markdown/文本"""

    @property
    def name(self) -> str:
        return "parse_document"

    @property
    def description(self) -> str:
        return (
            "解析各类文档文件（PDF/Word/PPT/Excel/Markdown/文本），"
            "提取其中的文本内容和表格数据。"
        )

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "要解析的文件路径",
                },
                "extract_tables": {
                    "type": "boolean",
                    "description": "是否提取表格数据",
                    "default": False,
                },
                "page_range": {
                    "type": "string",
                    "description": "页码范围（仅PDF），如 '1-5' 或 '1,3,5'，不指定则全部解析",
                },
            },
            "required": ["file_path"],
        }

    async def execute(
        self,
        file_path: str,
        extract_tables: bool = False,
        page_range: Optional[str] = None,
    ) -> str:
        """执行文档解析"""
        path = Path(file_path)
        if not path.exists():
            return f"错误: 文件不存在 - {file_path}"

        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                content = self._parse_pdf(path, page_range)
            elif suffix == ".docx":
                content = self._parse_docx(path)
            elif suffix == ".pptx":
                content = self._parse_pptx(path)
            elif suffix == ".xlsx":
                content = self._parse_xlsx(path)
            elif suffix in (".txt", ".md", ".py", ".json", ".yaml", ".yml"):
                content = self._parse_text(path)
            else:
                return f"错误: 不支持的文件格式 - {suffix}"

            return content
        except Exception as e:
            logger.error(f"文档解析失败 {file_path}: {e}")
            return f"错误: 解析失败 - {e}"

    def _parse_pdf(self, path: Path, page_range: Optional[str] = None) -> str:
        import fitz

        doc = fitz.open(str(path))
        pages_to_extract = self._parse_page_range(page_range, len(doc))

        text_parts = []
        for i in pages_to_extract:
            if 0 <= i < len(doc):
                text_parts.append(doc[i].get_text())
        doc.close()
        return "\n\n".join(text_parts)

    def _parse_docx(self, path: Path) -> str:
        from docx import Document

        doc = Document(str(path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                if row_text.strip():
                    parts.append(row_text)
        return "\n\n".join(parts)

    def _parse_pptx(self, path: Path) -> str:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- 幻灯片 {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                        if row_text.strip():
                            slide_texts.append(row_text)
            parts.append("\n".join(slide_texts))
        return "\n\n".join(parts)

    def _parse_xlsx(self, path: Path) -> str:
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_texts = [f"--- Sheet: {sheet_name} ---"]
            for row in ws.iter_rows(values_only=True):
                row_values = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_values):
                    sheet_texts.append(" | ".join(row_values))
            parts.append("\n".join(sheet_texts))
        wb.close()
        return "\n\n".join(parts)

    def _parse_text(self, path: Path) -> str:
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return path.read_text(encoding="gbk")

    def _parse_page_range(self, range_str: Optional[str], total: int) -> list:
        if not range_str:
            return list(range(total))
        pages = set()
        for part in range_str.split(","):
            part = part.strip()
            if "-" in part:
                start, end = part.split("-", 1)
                pages.update(range(int(start) - 1, int(end)))
            else:
                pages.add(int(part) - 1)
        return sorted([p for p in pages if 0 <= p < total])
